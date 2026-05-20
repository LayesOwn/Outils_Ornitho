# -*- coding: utf-8 -*-
"""
Génère les supports de formation ORNI-LAB en PDF et PowerPoint.
Usage : python generer_formation.py
Sortie : dossier docs/formation/export/
"""

import os
from pathlib import Path

# ── PowerPoint ─────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PDF (fpdf2) ─────────────────────────────────────────────────────────────
from fpdf import FPDF

# ─── Couleurs ORNI-LAB ──────────────────────────────────────────────────────
VERT_FONCE = RGBColor(0x1B, 0x43, 0x32)
VERT_MOYEN = RGBColor(0x2D, 0x6A, 0x4F)
VERT_CLAIR = RGBColor(0x52, 0xB7, 0x88)
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)

PDF_VERT_F = (27,  67,  50)
PDF_VERT   = (45, 106,  79)
PDF_VERT_C = (82, 183, 136)
PDF_GRIS   = (240, 247, 244)
PDF_GRIS2  = (224, 237, 232)
PDF_JAUNE  = (255, 248, 231)
PDF_BLANC  = (255, 255, 255)
PDF_NOIR   = (26,  26,  26)

cm = 10  # 1 cm = 10 mm

FONT_DIR = "C:/Windows/Fonts"

OUT_DIR = Path(__file__).parent / "export"
OUT_DIR.mkdir(exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ══════════════════════════════════════════════════════════════════════════

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _txb(slide, text, left, top, width, height,
         size=18, bold=False, color=BLANC, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _slide_titre(prs, titre, sous_titre):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, VERT_FONCE)
    _add_rect(sl, 0, 5.8, 10, 0.8, VERT_CLAIR)
    _txb(sl, titre, 0.4, 1.2, 9.2, 2.2, size=30, bold=True,
         color=BLANC, align=PP_ALIGN.CENTER)
    _txb(sl, sous_titre, 0.4, 3.6, 9.2, 0.8, size=17,
         color=VERT_CLAIR, align=PP_ALIGN.CENTER)
    _txb(sl, "ORNI-LAB  -  Abdoulaye Diop  -  2026",
         0.4, 5.85, 9.2, 0.5, size=11, color=BLANC, align=PP_ALIGN.CENTER)
    return sl


def _slide_contenu(prs, titre, puces, note_oral=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, BLANC)
    _add_rect(sl, 0, 0, 0.18, 7.5, VERT_FONCE)
    _add_rect(sl, 0.18, 0, 9.82, 1.1, VERT_MOYEN)
    _txb(sl, titre, 0.3, 0.08, 9.4, 0.9, size=22, bold=True,
         color=BLANC, align=PP_ALIGN.LEFT)
    txb = sl.shapes.add_textbox(Inches(0.45), Inches(1.25),
                                 Inches(9.3), Inches(5.5))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, puce in enumerate(puces):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = puce
        run.font.size = Pt(16)
        run.font.color.rgb = VERT_FONCE
    if note_oral:
        _txb(sl, "  " + note_oral, 0.25, 6.5, 9.5, 0.7,
             size=11, color=VERT_MOYEN, align=PP_ALIGN.LEFT)
    n = len(prs.slides)
    _txb(sl, str(n), 9.5, 6.8, 0.4, 0.35, size=10,
         color=VERT_CLAIR, align=PP_ALIGN.RIGHT)
    return sl


def _slide_2col(prs, titre, col_g, col_d, titre_g="", titre_d=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, BLANC)
    _add_rect(sl, 0, 0, 0.18, 7.5, VERT_FONCE)
    _add_rect(sl, 0.18, 0, 9.82, 1.1, VERT_MOYEN)
    _txb(sl, titre, 0.3, 0.08, 9.4, 0.9, size=22, bold=True,
         color=BLANC, align=PP_ALIGN.LEFT)
    _add_rect(sl, 5.05, 1.2, 0.05, 5.8, VERT_CLAIR)
    for col_x, items, col_titre in [
            (0.35, col_g, titre_g), (5.2, col_d, titre_d)]:
        if col_titre:
            _txb(sl, col_titre, col_x, 1.2, 4.5, 0.5,
                 size=14, bold=True, color=VERT_FONCE)
        txb = sl.shapes.add_textbox(Inches(col_x), Inches(1.75),
                                     Inches(4.5), Inches(5.0))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = item
            run.font.size = Pt(15)
            run.font.color.rgb = VERT_FONCE
    n = len(prs.slides)
    _txb(sl, str(n), 9.5, 6.8, 0.4, 0.35, size=10,
         color=VERT_CLAIR, align=PP_ALIGN.RIGHT)
    return sl


def generer_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_titre(prs,
        "ORNI-LAB\nLaboratoire pédagogique interactif\nen écologie quantitative des oiseaux",
        "Formation à l'outil  -  L3 / Master  -  [Date]  -  [Établissement]")

    _slide_contenu(prs, "Le problème pédagogique", [
        "- Les méthodes quantitatives en ornithologie sont abstraites à enseigner",
        "- R et Python demandent une courbe d'apprentissage importante",
        "- Les étudiants ont du mal à visualiser l'effet d'un paramètre sans coder",
        "",
        "=> ORNI-LAB permet d'explorer les modèles sans programmer,",
        "   de façon interactive, en temps réel.",
    ], note_oral="Insister : l'outil ne remplace pas R mais prépare à le comprendre.")

    _slide_2col(prs, "18 simulateurs interactifs",
        col_g=[
            "  Analyse CSV",
            "  Statistiques descriptives",
            "  Corrélation & régression",
            "  Tests statistiques",
            "  GLM Comptage",
            "  Modèles mixtes (LMM)",
            "  Domaine vital - MCP",
            "  Domaine vital - KDE",
        ],
        col_d=[
            "  Richesse & diversité",
            "  Croissance exponentielle / logistique",
            "  Matrices de Leslie",
            "  Capture-Marquage-Recapture",
            "  Modèles d'occupation",
            "  Distance sampling",
            "  Lotka-Volterra",
            "  Séries temporelles",
            "  PVA & conservation",
            "  Scénarios de gestion",
        ],
        titre_g="Biostatistique (8)",
        titre_d="Dynamique des populations (10)")

    _slide_2col(prs, "S'adapter à chaque public",
        col_g=[
            "MODE ÉTUDIANT",
            "",
            "- Interface épurée",
            "- Graphiques interactifs",
            "- Métriques clés",
            "- Interprétation guidée",
        ],
        col_d=[
            "MODE ENSEIGNANT",
            "",
            "- Tout le mode Étudiant +",
            "- Formules mathématiques (LaTeX)",
            "- Notes pédagogiques",
            "- Hypothèses implicites",
            "- Limites du modèle",
        ])

    _slide_contenu(prs, "Prise en main en 2 minutes", [
        "1.  Lancer ORNI-LAB.bat (ou : streamlit run app.py)",
        "2.  Observer la barre latérale : Mode -> Section -> Module",
        "3.  Choisir 'Croissance exponentielle / logistique'",
        "4.  Bouger le curseur  r  (taux de croissance)",
        "5.  Observer le changement en temps réel sur le graphique",
        "",
        "Survol (hover) : affiche les valeurs exactes",
        "Double-clic    : réinitialise le zoom",
    ], note_oral="Démo en direct. Montrer le hover et le zoom Plotly.")

    _slide_contenu(prs, "Travailler avec ses propres données", [
        "- Format CSV  -  séparateur auto-détecté (virgule, point-virgule, tabulation)",
        "- Colonnes numériques et catégoriques identifiées automatiquement",
        "- Valeurs manquantes détectées et signalées",
        "- Décimales : point  ou  virgule  acceptés",
        "",
        "Compatible avec :",
        "  -> Exports de tableurs (Excel, LibreOffice)",
        "  -> Loggers GPS (format terrain)",
        "  -> Applications de suivi (STOC, LPO, Faune-France...)",
    ])

    _slide_contenu(prs, "Module phare - Analyse de viabilité (PVA)", [
        "- ~250 trajectoires stochastiques simulées",
        "- Paramètres : N0,  r-moyen,  sigma (variabilité),  K,  seuil de quasi-extinction",
        "- Résultat : probabilité de quasi-extinction sur l'horizon choisi",
        "- Export PDF avec paramètres + résultat + interprétation",
        "",
        "Risque de quasi-extinction = proportion de simulations",
        "passant sous le seuil critique au moins une fois.",
        "",
        "=> Ce n'est pas une date d'extinction. C'est une mesure de fragilité.",
    ], note_oral="Un risque de 35% ne prédit pas l'extinction - il mesure la vulnérabilité sous les hypothèses choisies.")

    _slide_contenu(prs, "Module phare - Matrices de Leslie", [
        "- Entrée : fécondités et survies par classe d'âge",
        "- Calcul du taux de croissance asymptotique  lambda  (valeur propre dominante)",
        "",
        "   lambda > 1  =>  croissance à long terme",
        "   lambda = 1  =>  stabilité théorique",
        "   lambda < 1  =>  déclin à long terme",
        "",
        "- Sensibilité & élasticité : quel paramètre a le plus d'impact sur lambda ?",
        "- Application directe : prioriser les actions de conservation",
    ], note_oral="Ex. : élasticité survie adulte (0.52) >> fécondité (0.12) => protéger les adultes d'abord.")

    _slide_contenu(prs, "Comparer des actions de conservation", [
        "Trois scénarios comparés en parallèle :",
        "",
        "  A.  Restauration d'habitat       =>  augmenter K",
        "  B.  Réduction de la mortalité    =>  diminuer les pertes annuelles",
        "  C.  Translocation                =>  augmenter N0",
        "",
        "Critères de comparaison :",
        "  - Risque de quasi-extinction final",
        "  - Médiane de l'effectif projeté",
        "  - Faisabilité et coût (discussion)",
    ], note_oral="Les étudiants jouent le rôle de gestionnaires et justifient leur choix avec des chiffres.")

    _slide_contenu(prs, "Intégrer ORNI-LAB dans vos cours", [
        "Cours magistral   =>  démo live, visualiser l'effet d'un paramètre",
        "TP informatique   =>  fiches guidées par module (fournies)",
        "Travail personnel =>  exploration libre + export PDF/CSV",
        "Évaluation        =>  exporter un scénario, rédiger une note d'interprétation",
        "",
        "Cursus recommandé (4 x 3h) :",
        "  Séance 1 : Biostat de base    (CSV, stats, régression)",
        "  Séance 2 : Biostat avancées   (tests, GLM, LMM, GPS)",
        "  Séance 3 : Dynamique          (diversité, croissance, Leslie, CMR)",
        "  Séance 4 : Conservation       (occupation, PVA, scénarios)",
    ])

    _slide_contenu(prs, "Produire un rendu", [
        "PDF   =>  paramètres + résultats + interprétation  (compte-rendu, rapport)",
        "CSV   =>  données simulées ou analysées            (traitement R / Excel)",
        "PNG   =>  graphique via bouton Plotly              (présentation, rapport)",
        "",
        "Le PDF est généré en un clic depuis chaque module.",
        "Les symboles scientifiques (lambda, alpha, beta, sigma) sont normalisés automatiquement.",
    ])

    _slide_contenu(prs, "Ce que l'outil ne fait pas", [
        "- Les jeux de données par défaut sont SYNTHÉTIQUES",
        "  => Apprendre la méthode, pas produire une conclusion biologique",
        "",
        "- Une courbe bien ajustée  =/=  validation du mécanisme biologique",
        "- Une p-value faible       =/=  importance écologique",
        "- Un risque PVA dépend FORTEMENT des hypothèses sur sigma et le seuil",
        "- La détectabilité n'est pas corrigée dans tous les modules",
        "",
        "=> Le mode Enseignant liste explicitement les limites de chaque modèle.",
    ])

    _slide_contenu(prs, "Pour aller plus loin", [
        "docs/guide_scientifique.md",
        "  -> Variables, hypothèses, interprétations détaillées (18 modules)",
        "",
        "docs/formation/",
        "  -> Fiches pédagogiques biostat + dynamique",
        "  -> TP séances 1 à 4 (exercices guidés + questions)",
        "  -> Ce guide de présentation",
        "",
        "Mode Enseignant dans l'application",
        "  -> Notes pédagogiques et références dans chaque module",
    ])

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, VERT_FONCE)
    _add_rect(sl, 0, 5.8, 10, 0.8, VERT_CLAIR)
    _txb(sl, "Questions & Exploration libre", 0.4, 1.8, 9.2, 1.2,
         size=30, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    _txb(sl, "Choisissez un module, explorez, exportez un PDF.",
         0.4, 3.2, 9.2, 0.6, size=18, color=VERT_CLAIR, align=PP_ALIGN.CENTER)
    _txb(sl, "ORNI-LAB  -  Abdoulaye Diop  -  2026",
         0.4, 5.85, 9.2, 0.5, size=11, color=BLANC, align=PP_ALIGN.CENTER)

    out = OUT_DIR / "ORNI-LAB_Presentation.pptx"
    prs.save(str(out))
    print(f"  PPTX  ->  {out}")


# ══════════════════════════════════════════════════════════════════════════
#  PDF  (fpdf2)
# ══════════════════════════════════════════════════════════════════════════

class OrniPDF(FPDF):
    def __init__(self, doc_title="ORNI-LAB"):
        super().__init__('P', 'mm', 'A4')
        self.doc_title = doc_title
        self.set_auto_page_break(auto=True, margin=18)
        self.set_margins(18, 22, 18)
        self.add_font("Cal",  "",   f"{FONT_DIR}/calibri.ttf",  uni=True)
        self.add_font("Cal",  "B",  f"{FONT_DIR}/calibrib.ttf", uni=True)
        self.add_font("Cal",  "I",  f"{FONT_DIR}/calibrii.ttf", uni=True)
        self.add_font("Cal",  "BI", f"{FONT_DIR}/calibriz.ttf", uni=True)
        self.add_font("Mono", "",   f"{FONT_DIR}/cour.ttf",     uni=True)

    def header(self):
        self.set_fill_color(*PDF_VERT_F)
        self.rect(0, 0, 210, 10, 'F')
        self.set_text_color(*PDF_BLANC)
        self.set_font('Cal', 'B', 8)
        self.set_xy(0, 1.5)
        self.cell(0, 6, "  ORNI-LAB  -  Supports de formation", align='L')
        self.set_font('Cal', '', 8)
        self.set_xy(0, 1.5)
        self.cell(0, 6, f"Page {self.page_no()}  ", align='R')
        self.set_fill_color(*PDF_VERT_C)
        self.rect(0, 10, 210, 0.5, 'F')
        self.ln(6)

    def footer(self):
        self.set_y(-12)
        self.set_fill_color(*PDF_GRIS)
        self.rect(0, self.get_y(), 210, 12, 'F')
        self.set_text_color(*PDF_VERT)
        self.set_font('Cal', '', 7)
        self.cell(0, 8,
            "Abdoulaye Diop  -  Bioinformaticien / Biomathématicien  -  2026",
            align='C')

    def _bg_rect(self, color, h):
        self.set_fill_color(*color)
        self.rect(0, self.get_y(), 210, h, 'F')

    def bandeau_titre(self, titre, sous_titre=""):
        self._bg_rect(PDF_VERT_F, 16)
        self.set_text_color(*PDF_BLANC)
        self.set_font('Cal', 'B', 16)
        self.ln(4)
        self.multi_cell(self.epw, 9, titre, align='C')
        if sous_titre:
            self._bg_rect(PDF_VERT, 8)
            self.set_font('Cal', '', 9)
            self.multi_cell(self.epw, 7, sous_titre, align='C')
        self.set_fill_color(*PDF_VERT_C)
        self.rect(self.l_margin, self.get_y(), self.epw, 0.8, 'F')
        self.ln(5)
        self.set_text_color(*PDF_NOIR)

    def h1(self, texte):
        self.ln(3)
        self._bg_rect(PDF_VERT_F, 9)
        self.set_text_color(*PDF_BLANC)
        self.set_font('Cal', 'B', 11)
        self.set_x(self.l_margin)
        self.multi_cell(self.epw, 9, "  " + texte, align='L')
        self.ln(2)
        self.set_text_color(*PDF_NOIR)

    def h2(self, texte):
        self.ln(2)
        self.set_text_color(*PDF_VERT)
        self.set_font('Cal', 'B', 10)
        self.multi_cell(self.epw, 6, texte, align='L')
        self.set_fill_color(*PDF_VERT_C)
        self.rect(self.l_margin, self.get_y(), self.epw * 0.4, 0.4, 'F')
        self.ln(2)
        self.set_text_color(*PDF_NOIR)

    def h3(self, texte):
        self.ln(1)
        self.set_text_color(*PDF_VERT)
        self.set_font('Cal', 'BI', 9.5)
        self.multi_cell(self.epw, 5.5, texte, align='L')
        self.ln(1)
        self.set_text_color(*PDF_NOIR)

    def corps(self, texte):
        self.set_font('Cal', '', 9.5)
        self.set_text_color(*PDF_NOIR)
        self.multi_cell(self.epw, 5.5, texte, align='J')
        self.ln(1)

    def puce(self, texte, indent=0):
        self.set_font('Cal', '', 9.5)
        self.set_text_color(*PDF_NOIR)
        self.set_x(self.l_margin + indent)
        self.multi_cell(self.epw - indent, 5.5, texte, align='L')

    def puce2(self, texte):
        self.puce(texte, indent=8)

    def code(self, texte):
        self.ln(1)
        self.set_fill_color(*PDF_GRIS2)
        self.set_text_color(*PDF_VERT_F)
        self.set_font('Mono', '', 8.5)
        self.multi_cell(self.epw, 5, texte, fill=True, align='L')
        self.ln(2)
        self.set_text_color(*PDF_NOIR)

    def note(self, texte):
        self.ln(1)
        self.set_fill_color(*PDF_JAUNE)
        self.set_text_color(100, 80, 0)
        self.set_font('Cal', 'I', 9)
        self.multi_cell(self.epw, 5.5, texte, fill=True, align='L')
        self.ln(2)
        self.set_text_color(*PDF_NOIR)

    def sep(self):
        self.ln(3)
        self.set_draw_color(*PDF_VERT_C)
        self.set_line_width(0.4)
        self.line(self.l_margin, self.get_y(),
                  self.l_margin + self.epw, self.get_y())
        self.ln(4)

    def tableau(self, data, col_widths=None):
        self.ln(2)
        if col_widths is None:
            n = len(data[0])
            col_widths = [self.epw / n] * n

        for i, row in enumerate(data):
            is_header = (i == 0)
            if is_header:
                self.set_fill_color(*PDF_VERT)
                self.set_text_color(*PDF_BLANC)
                self.set_font('Cal', 'B', 8.5)
            else:
                self.set_fill_color(*(PDF_BLANC if i % 2 == 1 else PDF_GRIS))
                self.set_text_color(*PDF_NOIR)
                self.set_font('Cal', '', 8.5)

            x0 = self.l_margin
            y0 = self.get_y()

            max_h = 7
            for j, cell_txt in enumerate(row):
                self.set_font('Cal', 'B' if is_header else '', 8.5)
                nb = self.get_string_width(str(cell_txt))
                n_lines = max(1, int(nb / max(1, col_widths[j] - 4)) + 1)
                max_h = max(max_h, n_lines * 5.5 + 3)

            for j, cell_txt in enumerate(row):
                self.set_xy(x0, y0)
                fill_color = PDF_VERT if is_header else (
                    PDF_BLANC if i % 2 == 1 else PDF_GRIS)
                self.set_fill_color(*fill_color)
                self.rect(x0, y0, col_widths[j], max_h, 'F')
                self.set_draw_color(*PDF_VERT_C)
                self.set_line_width(0.3)
                self.rect(x0, y0, col_widths[j], max_h)
                self.set_xy(x0 + 2, y0 + 1.5)
                self.set_font('Cal', 'B' if is_header else '', 8.5)
                txt_color = PDF_BLANC if is_header else PDF_NOIR
                self.set_text_color(*txt_color)
                self.multi_cell(col_widths[j] - 4, 5.5,
                                str(cell_txt), align='L')
                x0 += col_widths[j]

            self.set_y(y0 + max_h)

        self.ln(4)
        self.set_draw_color(0, 0, 0)
        self.set_text_color(*PDF_NOIR)


def _build_pdf(nom_fichier, titre, sous_titre, sections):
    """
    sections : liste de tuples (type, contenu)
    types supportés : h1, h2, h3, corps, puce, puce2, code, note, sep, saut,
                      tableau(data, col_widths)
    """
    pdf = OrniPDF(titre)
    pdf.add_page()
    pdf.bandeau_titre(titre, sous_titre)

    for item in sections:
        kind = item[0]
        if kind == "saut":
            pdf.add_page()
        elif kind == "sep":
            pdf.sep()
        elif kind == "tableau":
            _, data, widths = item
            pdf.tableau(data, widths)
        else:
            _, texte = item
            getattr(pdf, kind)(texte)

    out = OUT_DIR / nom_fichier
    pdf.output(str(out))
    print(f"  PDF   ->  {out}")


# ══════════════════════════════════════════════════════════════════════════
#  CONTENU DES PDF
# ══════════════════════════════════════════════════════════════════════════

def pdf_programme():
    _build_pdf("01_Programme_formation.pdf",
        "ORNI-LAB - Programme de formation",
        "Cursus 4 séances x 3h  -  L3 / Master  -  Écologie quantitative des oiseaux",
        [
        ("h1", "Présentation générale"),
        ("corps",
            "ORNI-LAB est un laboratoire pédagogique interactif pour l'écologie quantitative "
            "et la modélisation des populations d'oiseaux. Il regroupe 18 simulateurs organisés "
            "en deux sections : Biostatistique et Dynamique des populations."),
        ("corps",
            "L'outil est conçu pour des étudiants de niveau L3 à Master, des enseignants "
            "souhaitant l'intégrer dans leurs cours, et des professionnels de terrain."),

        ("h1", "Les deux modes d'utilisation"),
        ("puce", "- Mode Étudiant : interface épurée, interprétations guidées, focus sur les résultats"),
        ("puce", "- Mode Enseignant : formules LaTeX, notes pédagogiques, hypothèses et limites de chaque modèle"),

        ("h1", "Cursus recommandé - 4 séances de 3 heures"),
        ("tableau",
            [["Séance", "Thème", "Modules utilisés"],
             ["1  (3h)", "Prise en main + Biostat de base",
              "Analyse CSV, Statistiques descriptives, Corrélation & régression"],
             ["2  (3h)", "Biostatistiques avancées",
              "Tests statistiques, GLM Comptage, Modèles mixtes, Domaines vitaux MCP/KDE"],
             ["3  (3h)", "Dynamique des populations - fondamentaux",
              "Richesse & diversité, Croissance, Matrices de Leslie, CMR"],
             ["4  (3h)", "Dynamique avancée & Conservation",
              "Occupation, Distance sampling, Lotka-Volterra, Séries temporelles, PVA, Scénarios"],
            ],
            [1.5*cm, 5.0*cm, 9.0*cm]),

        ("h1", "Ressources disponibles"),
        ("tableau",
            [["Fichier", "Description"],
             ["01_GUIDE_INSTALLATION.md", "Installation, interface, import CSV, exports"],
             ["03_PRESENTATION_DIAPORAMA.md", "Plan de présentation 45 min (14 diapositives)"],
             ["fiches/FICHES_BIOSTATISTIQUE.md", "8 fiches pédagogiques - section Biostatistique"],
             ["fiches/FICHES_DYNAMIQUE.md", "10 fiches pédagogiques - section Dynamique"],
             ["TP/TP_SEANCE1 à TP_SEANCE4.md", "Exercices guidés par séance avec questions"],
             ["../guide_scientifique.md", "Référence scientifique : variables, hypothèses, interprétations"],
            ],
            [6.5*cm, 9.0*cm]),

        ("h1", "Conseils pédagogiques généraux"),
        ("puce", "1. Demandez toujours aux étudiants de prédire le résultat avant de bouger un curseur."),
        ("puce", "2. Modifiez une variable à la fois pour isoler son effet."),
        ("puce", "3. Comparez toujours au moins deux scénarios avant de conclure."),
        ("puce", "4. Distinguez : effet statistique =/= effet biologique =/= effet de gestion."),
        ("puce", "5. Notez les hypothèses non représentées dans le simulateur."),
        ("puce", "6. Utilisez les exports CSV pour faire recalculer un résultat hors ORNI-LAB."),
        ])


def pdf_guide_installation():
    _build_pdf("02_Guide_installation.pdf",
        "Guide d'installation et prise en main",
        "ORNI-LAB  -  Lancement, interface, import de données, exports",
        [
        ("h1", "1. Prérequis"),
        ("tableau",
            [["Logiciel", "Version minimale", "Vérification"],
             ["Python", "3.10+", "python --version"],
             ["pip", "inclus avec Python", "pip --version"],
             ["Navigateur", "Chrome / Firefox / Edge", "-"],
            ],
            [4.5*cm, 5.0*cm, 6.0*cm]),
        ("note", "Python peut être téléchargé sur python.org. Cocher 'Add Python to PATH' lors de l'installation."),

        ("h1", "2. Installation"),
        ("h2", "Option A - Lancement automatique Windows (recommandé)"),
        ("puce", "Double-cliquer sur ORNI-LAB.bat à la racine du projet."),
        ("puce", "Le fichier bat détecte Python, vérifie les dépendances et ouvre l'application."),
        ("h2", "Option B - Lancement manuel (PowerShell)"),
        ("code",
         "cd 'C:\\...\\Outils_Ornitho'\n"
         "python -m venv .venv\n"
         ".venv\\Scripts\\activate\n"
         "pip install -r requirements.txt\n"
         "python -m streamlit run app.py"),
        ("corps", "L'application s'ouvre automatiquement sur http://localhost:8501"),

        ("h1", "3. Interface - Vue d'ensemble"),
        ("tableau",
            [["Élément", "Description"],
             ["Barre latérale", "Choix du mode, de la section, du module et des paramètres (sliders)"],
             ["Zone principale", "Graphique interactif, métriques, interprétation guidée"],
             ["Boutons export", "Télécharger les données (CSV) ou le résumé (PDF)"],
            ],
            [5.0*cm, 10.5*cm]),

        ("h1", "4. Les deux modes"),
        ("h2", "Mode Étudiant"),
        ("puce", "- Interface épurée  -  Graphiques et métriques clés  -  Interprétation guidée"),
        ("puce", "- Recommandé pour les TP"),
        ("h2", "Mode Enseignant"),
        ("puce", "- Tout ce qui est en Mode Étudiant, plus :"),
        ("puce2", "-> Formules mathématiques (LaTeX)"),
        ("puce2", "-> Notes pédagogiques et discussion"),
        ("puce2", "-> Hypothèses implicites du modèle"),
        ("puce2", "-> Limites de chaque approche"),
        ("puce", "- Recommandé pour préparer un cours ou animer une séance"),

        ("h1", "5. Importer ses données CSV"),
        ("puce", "- Séparateurs acceptés : virgule ,   point-virgule ;   tabulation \\t"),
        ("puce", "- Décimales : point  .  ou virgule  ,"),
        ("puce", "- Encodage : UTF-8, Latin-1"),
        ("h2", "Procédure"),
        ("puce", "1. Dans la barre latérale, sélectionner 'Analyse CSV'"),
        ("puce", "2. Cliquer sur 'Parcourir' et choisir le fichier .csv"),
        ("puce", "3. Le module détecte automatiquement le séparateur et l'encodage"),
        ("puce", "4. Les colonnes numériques et catégoriques sont identifiées automatiquement"),

        ("h1", "6. Utiliser un graphique Plotly"),
        ("tableau",
            [["Action", "Effet"],
             ["Clic + glisser", "Zoomer sur une zone"],
             ["Double-clic", "Réinitialiser le zoom"],
             ["Survol (hover)", "Afficher les valeurs exactes"],
             ["Clic sur légende", "Masquer / afficher une série"],
             ["Icônes haut-droite", "Télécharger le graphique en PNG"],
            ],
            [5.0*cm, 10.5*cm]),

        ("h1", "7. Exports"),
        ("puce", "- CSV : données simulées ou analysées, compatible Excel / R / Python"),
        ("puce", "- PDF : résumé A4 (paramètres + résultats + interprétation)"),

        ("h1", "8. Résolution de problèmes fréquents"),
        ("tableau",
            [["Problème", "Solution"],
             ["L'app ne se lance pas", "Vérifier que Python est dans le PATH"],
             ["Page blanche", "Attendre 10 s ou rafraîchir (F5)"],
             ["ModuleNotFoundError", "Relancer : pip install -r requirements.txt"],
             ["CSV non reconnu", "Vérifier le séparateur et l'encodage"],
             ["Port 8501 occupé", "streamlit run app.py --server.port 8502"],
            ],
            [5.5*cm, 10.0*cm]),
        ])


def pdf_fiches_biostat():
    _build_pdf("03_Fiches_Biostatistique.pdf",
        "Fiches pédagogiques - Biostatistique",
        "8 modules  -  Objectif, concept clé, manipulation guidée, questions d'interprétation",
        [
        ("h1", "Fiche 1 - Analyse CSV"),
        ("h2", "Objectif pédagogique"),
        ("corps", "Explorer et nettoyer un fichier de terrain réel. Identifier les problèmes "
                  "de qualité de données avant toute analyse."),
        ("h2", "Concept clé"),
        ("corps", "Un fichier de terrain brut contient presque toujours des valeurs manquantes, "
                  "des unités mélangées ou des colonnes mal typées. L'exploration préalable est "
                  "une étape incontournable."),
        ("h2", "Manipulations guidées"),
        ("puce", "1. Importer le fichier CSV (séparateur auto-détecté)"),
        ("puce", "2. Lire le résumé automatique : lignes, colonnes, valeurs manquantes"),
        ("puce", "3. Sélectionner une colonne numérique => observer histogramme + boîte à moustaches"),
        ("puce", "4. Choisir deux variables => tracer le nuage de points et lire R²"),
        ("puce", "5. Comparer deux groupes avec le test de Welch => lire p-value"),
        ("puce", "6. Exporter CSV nettoyé + PDF résumé"),
        ("h2", "Questions d'interprétation"),
        ("puce", "Q1. La moyenne est supérieure à la médiane pour 'comptage'. Qu'est-ce que cela indique ?"),
        ("puce", "Q2. CV = 85 %. Que dit cela sur l'hétérogénéité inter-sites ?"),
        ("puce", "Q3. p-value = 0.23 (Welch). Peut-on conclure que les habitats sont identiques ?"),
        ("puce", "Q4. r = 0.91 entre deux colonnes. Est-ce une preuve de causalité ?"),
        ("h2", "Erreurs fréquentes"),
        ("puce", "X  Mélanger des unités dans une même colonne (g et kg)"),
        ("puce", "X  Interpréter corrélation comme causalité"),
        ("puce", "X  Ignorer les valeurs manquantes avant de calculer une moyenne"),
        ("sep", ""),

        ("h1", "Fiche 2 - Statistiques descriptives"),
        ("h2", "Objectif pédagogique"),
        ("corps", "Résumer une distribution de comptages ou de mesures morphologiques. "
                  "Choisir le bon indicateur selon la forme de la distribution."),
        ("h2", "Concept clé"),
        ("corps", "La moyenne seule ne suffit pas. Une distribution asymétrique (fréquente en "
                  "ornithologie) nécessite la médiane et les quartiles. L'écart-type mesure la "
                  "dispersion, pas la forme."),
        ("h2", "Manipulations guidées"),
        ("puce", "1. Régler nombre de sites, abondance moyenne, agrégation spatiale"),
        ("puce", "2. Comparer moyenne vs médiane selon le niveau d'agrégation"),
        ("puce", "3. Augmenter le nombre de sites => observer la stabilisation des indicateurs"),
        ("h2", "Questions d'interprétation"),
        ("puce", "Q1. Avec forte agrégation : moyenne = 18, médiane = 8. Quelle valeur représente 70% des sites ?"),
        ("puce", "Q2. L'écart-type est plus grand que la moyenne. Qu'implique-t-on pour un test de normalité ?"),
        ("puce", "Q3. Pourquoi les ornithologues utilisent-ils des modèles de comptage (Poisson) ?"),
        ("sep", ""),

        ("h1", "Fiche 3 - Corrélation et régression"),
        ("h2", "Objectif pédagogique"),
        ("corps", "Quantifier la relation linéaire entre deux variables biologiques. Distinguer la "
                  "force d'association (r, R²) de la signification statistique (p-value)."),
        ("h2", "Formules"),
        ("code", "r de Pearson  =  Cov(X,Y) / (sigma_X x sigma_Y)   dans [-1 ; 1]\n"
                 "R²            =  r²   (proportion de variance expliquée)\n"
                 "Y_prédit      =  a + b x X"),
        ("h2", "Questions d'interprétation"),
        ("puce", "Q1. n = 200, r = 0.15, p = 0.02. La relation est-elle biologiquement importante ?"),
        ("puce", "Q2. Pente = 2.3 g/mm. Que signifie ce chiffre pour masse ~ longueur d'aile ?"),
        ("puce", "Q3. R² = 0.62. Quelle proportion de la variance n'est pas expliquée ?"),
        ("h2", "Erreurs fréquentes"),
        ("puce", "X  Confondre r (corrélation) et R² (proportion de variance)"),
        ("puce", "X  P-value faible =/= relation forte  |  R² élevé =/= causalité"),
        ("puce", "X  Interpréter l'intercept quand X = 0 est hors domaine biologique"),
        ("sep", ""),

        ("h1", "Fiche 4 - Tests statistiques"),
        ("h2", "Objectif pédagogique"),
        ("corps", "Comparer des groupes biologiques (habitats, années, traitements). "
                  "Choisir le bon test selon la structure des données."),
        ("tableau",
            [["Test", "Quand l'utiliser"],
             ["t-test de Welch", "2 groupes, variable continue, variances potentiellement inégales"],
             ["ANOVA", "3 groupes ou plus, variable continue"],
             ["Mann-Whitney", "Alternative non-paramétrique au t-test"],
             ["Khi-deux", "Comparer des fréquences ou proportions"],
            ],
            [4.5*cm, 11.0*cm]),
        ("h2", "Questions d'interprétation"),
        ("puce", "Q1. p = 0.04 avec n = 10. Que faire avant de conclure ?"),
        ("puce", "Q2. Deux habitats : moyennes 12 et 18, p = 0.001. La différence est-elle biologiquement importante ?"),
        ("puce", "Q3. Interpréter p > 0.05 comme 'les groupes sont identiques' - est-ce correct ?"),
        ("sep", ""),

        ("h1", "Fiche 5 - GLM Comptage (Poisson / Binomial négatif)"),
        ("h2", "Concept clé"),
        ("corps", "Les comptages (0, 1, 2, ...) ne suivent pas une loi normale. Le GLM Poisson "
                  "est la base ; le binomial négatif gère la surdispersion (variance > moyenne) "
                  "fréquente en ornithologie."),
        ("h2", "Sélection du modèle"),
        ("puce", "1. Ajuster GLM Poisson => observer ratio Déviance/ddl"),
        ("puce", "2. Ratio >> 1 => surdispersion => basculer sur Binomial négatif"),
        ("puce", "3. Comparer les AIC (plus bas = mieux)"),
        ("code", "exp(beta)  =  effet multiplicatif d'une covariable sur le comptage moyen\n"
                 "exp(0.85) = 2.34  =>  2.34 fois plus d'individus dans le groupe de référence"),
        ("sep", ""),

        ("h1", "Fiche 6 - Modèles mixtes (LMM)"),
        ("h2", "Concept clé"),
        ("corps", "Quand plusieurs mesures viennent du même site ou individu, elles ne sont pas "
                  "indépendantes (pseudo-réplication). Un modèle mixte ajoute des effets aléatoires "
                  "pour capturer cette structure."),
        ("tableau",
            [["Terme", "Définition"],
             ["Effet fixe", "Variable d'intérêt (habitat, traitement) - estimée en moyenne"],
             ["Effet aléatoire", "Variable de groupement (site, individu) - variabilité contrôlée"],
             ["ICC", "Corrélation intra-classe : proportion de variance due au groupement"],
            ],
            [4.0*cm, 11.5*cm]),
        ("puce", "Q1. ICC = 0.05 : le modèle mixte est-il nécessaire ?"),
        ("puce", "Q2. Données : 10 obs/oiseau bagué sur 5 individus => effet aléatoire recommandé ?"),
        ("sep", ""),

        ("h1", "Fiche 7 - Domaines vitaux : MCP et KDE"),
        ("tableau",
            [["Méthode", "Description", "Avantage", "Limite principale"],
             ["MCP", "Polygone convexe minimum contenant tous les points",
              "Simple, universel", "Sensible aux points extrêmes"],
             ["KDE", "Distribution de probabilité de présence (noyau gaussien)",
              "Probabiliste, isopleths ajustables", "Paramètre h à choisir, plus de points requis"],
            ],
            [2.0*cm, 5.0*cm, 4.0*cm, 4.5*cm]),
        ("h2", "Paramètres clés"),
        ("puce", "- MCP 95% : exclure 5% des points aberrants => surface réduite"),
        ("puce", "- KDE 50% : zone noyau d'utilisation intensive"),
        ("puce", "- KDE 95% : domaine vital complet"),
        ("puce", "- Bande passante h : fort lissage => contours lisses ; faible => sur-ajustement"),
        ("h2", "Questions d'interprétation"),
        ("puce", "Q1. MCP 100% = 1850 ha, MCP 95% = 620 ha. Que révèle cet écart ?"),
        ("puce", "Q2. L'isoplèthe KDE 50% représente quoi biologiquement ?"),
        ("puce", "Q3. Pourquoi ne pas comparer des MCP calculés avec des n différents ?"),

        ("h1", "Récapitulatif - Choisir la bonne méthode"),
        ("tableau",
            [["Situation", "Module recommandé"],
             ["Résumer des données de terrain", "Analyse CSV + Statistiques descriptives"],
             ["Deux variables continues", "Corrélation & régression"],
             ["Deux groupes à comparer", "Tests statistiques (t-test ou Mann-Whitney)"],
             ["Trois groupes ou plus", "Tests statistiques (ANOVA)"],
             ["Modéliser des comptages", "GLM Comptage"],
             ["Données groupées par site/individu", "Modèles mixtes"],
             ["Positions GPS => domaine vital géométrique", "Domaine vital MCP"],
             ["Positions GPS => domaine vital probabiliste", "Domaine vital KDE"],
            ],
            [7.5*cm, 8.0*cm]),
        ])


def pdf_fiches_dynamique():
    _build_pdf("04_Fiches_Dynamique_populations.pdf",
        "Fiches pédagogiques - Dynamique des populations",
        "10 modules  -  Objectif, concept clé, manipulation guidée, questions d'interprétation",
        [
        ("h1", "Fiche 9 - Richesse & Diversité"),
        ("tableau",
            [["Indice", "Ce qu'il mesure", "Sensible aux espèces rares ?"],
             ["Richesse S", "Nombre d'espèces", "Oui (chaque espèce compte pareil)"],
             ["Shannon H'", "Diversité + équitabilité", "Modérément"],
             ["Simpson 1-D", "Dominance inversée", "Non (pondéré par abondance)"],
             ["Équitabilité J", "Régularité des abondances", "-"],
            ],
            [3.5*cm, 6.5*cm, 5.5*cm]),
        ("h2", "Questions clés"),
        ("puce", "Q1. Un boisement a S = 18, une prairie a S = 15 mais H' plus élevé. Laquelle est plus diverse ?"),
        ("puce", "Q2. La courbe d'accumulation n'atteint pas l'asymptote. Que faire ?"),
        ("sep", ""),

        ("h1", "Fiche 10 - Croissance exponentielle et logistique"),
        ("tableau",
            [["Modèle", "Équation", "Hypothèse"],
             ["Exponentiel", "dN/dt = rN", "Pas de limite de ressources"],
             ["Logistique", "dN/dt = rN(1 - N/K)", "Ralentissement quand N => K"],
            ],
            [3.5*cm, 5.0*cm, 7.0*cm]),
        ("puce", "Q1. r = -0.05. Que se passe-t-il à long terme ?"),
        ("puce", "Q2. Doubler K double-t-il la population à l'an 20 ?"),
        ("puce", "Q3. À quel N la croissance nette est-elle maximale ?  =>  N = K/2"),
        ("sep", ""),

        ("h1", "Fiche 11 - Matrices de Leslie"),
        ("corps", "La matrice de Leslie organise les transitions entre classes d'âge. "
                  "La valeur propre dominante lambda donne le taux de croissance asymptotique."),
        ("code", "lambda > 1  =>  croissance    |    lambda = 1  =>  stabilité    |    lambda < 1  =>  déclin"),
        ("h2", "Sensibilité et élasticité"),
        ("puce", "- Sensibilité : dérivée de lambda par rapport à chaque paramètre (valeur absolue)"),
        ("puce", "- Élasticité  : variation proportionnelle - comparable entre paramètres"),
        ("puce", "- Application directe : paramètre avec élasticité max = priorité conservation"),
        ("h2", "Questions clés"),
        ("puce", "Q1. lambda = 0.93. Combien d'années avant une réduction de moitié ? (T = ln2/|ln(lambda)|)"),
        ("puce", "Q2. Élasticité S_adulte = 0.48 vs F_adulte = 0.12. Quelle intervention prioriser ?"),
        ("sep", ""),

        ("h1", "Fiche 12 - Capture-Marquage-Recapture"),
        ("code", "N_estimé = (M x C) / R\nIC 95% dépend directement de R : plus R est faible, plus l'incertitude est grande."),
        ("tableau",
            [["Paramètre", "Description"],
             ["M", "Individus marqués lors de la 1re session"],
             ["C", "Individus capturés lors de la 2e session"],
             ["R", "Individus marqués recapturés (variable critique)"],
            ],
            [2.5*cm, 13.0*cm]),
        ("h2", "Hypothèses à respecter"),
        ("puce", "1. Population fermée entre les deux sessions"),
        ("puce", "2. Marques conservées et reconnues"),
        ("puce", "3. Probabilité de capture identique pour tous les individus"),
        ("puce", "4. Mélange homogène des individus marqués"),
        ("sep", ""),

        ("h1", "Fiche 13 - Modèles d'occupation"),
        ("code", "Occupation observée (naïve) = psi x p    (sous-estimation systématique)\n"
                 "Le modèle MacKenzie sépare psi et p à partir de visites répétées."),
        ("puce", "Q1. psi = 0.70, p = 0.35 => occupation naïve = 0.25. Biais = 64% !"),
        ("puce", "Q2. Minimum recommandé : 3 visites répétées par site"),
        ("sep", ""),

        ("h1", "Fiche 14 - Distance Sampling"),
        ("code", "D = n / (2 x L x w_bar)\n"
                 "n = individus détectés  |  L = longueur transects  |  w_bar = distance effective"),
        ("puce", "Fonctions de détection : half-normal (décroissance douce) ou hazard-rate (épaule)"),
        ("puce", "Sélection par AIC. CV < 20% = précision acceptable."),
        ("sep", ""),

        ("h1", "Fiche 15 - Lotka-Volterra"),
        ("code", "dV/dt = alpha*V - beta*V*P     (Proies)\n"
                 "dP/dt = delta*V*P - gamma*P    (Prédateurs)"),
        ("puce", "Les oscillations sont déphasées : proies en avance d'environ 1/4 de cycle."),
        ("puce", "Q1. Augmenter alpha => l'effectif moyen des prédateurs augmente ou diminue ?"),
        ("sep", ""),

        ("h1", "Fiche 16 - Séries temporelles"),
        ("tableau",
            [["Méthode", "Type", "Usage principal"],
             ["Mann-Kendall", "Non-paramétrique", "Tendance monotone - tau et p-value"],
             ["Pente de Sen", "Non-paramétrique", "Magnitude de la tendance (individus/an)"],
             ["Régression linéaire", "Paramétrique", "Tendance linéaire supposée"],
            ],
            [4.0*cm, 4.0*cm, 7.5*cm]),
        ("puce", "Q1. tau = -0.38, p = 0.01. Conclusion sur la tendance ?"),
        ("puce", "Q2. Pente Sen = -3.2 ind/an, N0 = 85. En combien d'années N divisé par 2 ?"),
        ("sep", ""),

        ("h1", "Fiche 17 - PVA & Conservation"),
        ("tableau",
            [["Paramètre", "Rôle"],
             ["N0", "Effectif initial - plus faible = risque démographique plus élevé"],
             ["r-moyen", "Tendance moyenne - positif = croissance, mais variance peut dominer"],
             ["sigma (sd_r)", "Stochasticité environnementale - variable critique souvent sous-estimée"],
             ["K", "Capacité de charge (plafond lié à l'habitat)"],
             ["Seuil", "Effectif critique : décision arbitraire du modélisateur"],
             ["Pertes", "Mortalité additive annuelle (collisions, prélèvements, destructions)"],
            ],
            [3.0*cm, 12.5*cm]),
        ("note", "Risque = proportion de simulations passant sous le seuil. Ce n'est pas une date d'extinction."),
        ("sep", ""),

        ("h1", "Fiche 18 - Scénarios de gestion"),
        ("corps", "Comparer l'efficacité relative de différentes actions de conservation "
                  "sur une population menacée."),
        ("tableau",
            [["Action", "Paramètre modifié", "Exemple"],
             ["Restauration d'habitat", "Augmenter K", "Replantation de haies, zones tampon"],
             ["Réduction de mortalité", "Diminuer les pertes", "Dépose lignes, réduction collisions"],
             ["Translocation", "Augmenter N0", "Renforcement de population, réintroduction"],
            ],
            [4.5*cm, 5.0*cm, 6.0*cm]),
        ("note", "Critères de comparaison : risque de quasi-extinction final + médiane de l'effectif projeté + faisabilité terrain."),

        ("h1", "Récapitulatif - Choisir la bonne méthode"),
        ("tableau",
            [["Question biologique", "Module recommandé"],
             ["Comparer des communautés", "Richesse & diversité"],
             ["Projeter N à t+50 ans", "Croissance exponentielle/logistique"],
             ["Quelle classe d'âge protéger ?", "Matrices de Leslie"],
             ["Combien d'individus dans cette population ?", "CMR"],
             ["L'espèce occupe-t-elle ce site ?", "Modèles d'occupation"],
             ["Quelle densité sur mes transects ?", "Distance sampling"],
             ["Proie-prédateur : oscillations", "Lotka-Volterra"],
             ["Tendance sur 20 ans de comptages", "Séries temporelles"],
             ["Risque d'extinction à 50 ans ?", "PVA & Conservation"],
             ["Quelle action de conservation prioriser ?", "Scénarios de gestion"],
            ],
            [8.0*cm, 7.5*cm]),
        ])


def pdf_tp1():
    _build_pdf("05_TP_Seance1.pdf",
        "TP Séance 1 - Prise en main et biostatistiques de base",
        "Durée 3h  -  Prérequis : aucun  -  Modules : Analyse CSV, Stats desc., Régression",
        [
        ("h1", "Objectifs"),
        ("puce", "- Installer et lancer ORNI-LAB"),
        ("puce", "- Importer et explorer un fichier CSV de terrain"),
        ("puce", "- Calculer des statistiques descriptives et interpréter une régression linéaire"),

        ("h1", "Partie 0 - Installation (30 min)"),
        ("puce", "1. Double-cliquer sur ORNI-LAB.bat"),
        ("puce", "2. L'application s'ouvre sur http://localhost:8501"),
        ("puce", "3. Basculer entre Mode Étudiant et Mode Enseignant"),
        ("note", "Q0.1  Citer deux informations supplémentaires visibles en Mode Enseignant."),

        ("h1", "Partie 1 - Analyse CSV (45 min)"),
        ("corps", "Fichier exemple (colonnes : site, habitat, espèce, comptage, année)."),
        ("puce", "1. Importer le fichier CSV dans le module Analyse CSV"),
        ("puce", "2. Observer le résumé automatique"),
        ("puce", "3. Sélectionner la colonne comptage => histogramme + boîte à moustaches"),
        ("puce", "4. Comparer deux habitats avec le test de Welch"),
        ("puce", "5. Exporter le résumé en PDF"),
        ("note", "Q1.1  La distribution est-elle symétrique ?\n"
                 "Q1.2  p-value Welch = 0.23. Peut-on conclure que les habitats sont identiques ?\n"
                 "Q1.3  La différence est-elle biologiquement importante ?"),

        ("h1", "Partie 2 - Statistiques descriptives (45 min)"),
        ("puce", "1. Régler : 60 sites, abondance moyenne = 10, agrégation faible"),
        ("puce", "2. Lire moyenne, médiane, écart-type"),
        ("puce", "3. Augmenter l'agrégation => observer l'écart moyenne/médiane"),
        ("puce", "4. Réduire le nombre de sites à 10, puis remonter à 200"),
        ("note", "Q2.1  Avec forte agrégation : moyenne = 18, médiane = 8. Quelle valeur représente 70% des sites ?\n"
                 "Q2.2  Quel est l'effet de l'augmentation du nombre de sites sur la précision des indicateurs ?"),

        ("h1", "Partie 3 - Corrélation et régression (45 min)"),
        ("corps", "Contexte : relation entre longueur d'aile (mm) et masse (g) chez la Fauvette des jardins."),
        ("puce", "1. n = 15, pente = 3.5, variabilité faible => lire r, R², p-value"),
        ("puce", "2. Passer à n = 150 => observer l'évolution de la p-value"),
        ("puce", "3. Régler pente = 0 => que se passe-t-il ?"),
        ("note", "Q3.1  n = 15, r = 0.55 => p-value significative ?\n"
                 "Q3.2  La relation est-elle plus 'forte' avec n = 150 ? Justifier.\n"
                 "Q3.3  Pente = 2.8 g/mm, intercept = 5.2. Prédire la masse pour une aile de 70 mm.\n"
                 "Q3.4  Pourquoi l'intercept est-il difficile à interpréter si X = 0 est hors domaine biologique ?"),

        ("h1", "Récapitulatif des formules"),
        ("tableau",
            [["Indicateur", "Formule / Définition"],
             ["Médiane", "Valeur centrale de la distribution ordonnée"],
             ["Écart-type", "racine(Somme(xi - x_barre)² / (n-1))"],
             ["r de Pearson", "Cov(X,Y) / (sigma_X x sigma_Y)  dans [-1 ; 1]"],
             ["R²", "r² - proportion de variance expliquée"],
             ["p-value", "P(résultat aussi extrême | H0 vraie)"],
            ],
            [5.0*cm, 10.5*cm]),
        ])


def pdf_tp2():
    _build_pdf("06_TP_Seance2.pdf",
        "TP Séance 2 - Biostatistiques avancées",
        "Durée 3h  -  Prérequis : Séance 1  -  Modules : Tests, GLM, LMM, Domaines vitaux",
        [
        ("h1", "Partie 1 - Tests statistiques (45 min)"),
        ("corps", "Contexte : succès reproducteur de la Chouette hulotte en forêt feuillue vs résineuse."),
        ("puce", "1. t-test Welch : n = 25 par groupe, moyennes 2.8 vs 2.1 jeunes"),
        ("puce", "2. Réduire à n = 8 par groupe => observer la p-value"),
        ("puce", "3. Distribution asymétrique => comparer t-test vs Mann-Whitney"),
        ("puce", "4. ANOVA sur 3 types de forêts"),
        ("note", "Q1.1  La différence est significative avec n = 25. Quelle est la taille de l'effet ?\n"
                 "Q1.2  La même différence biologique devient non-significative avec n = 8. Que cela illustre-t-il ?\n"
                 "Q1.3  Dans quel cas préférer Mann-Whitney ?"),

        ("h1", "Partie 2 - GLM Comptage (45 min)"),
        ("corps", "Contexte : abondance du Bruant proyer sur 60 placettes - prairie intensive vs extensive."),
        ("puce", "1. Simuler comptages => ajuster GLM Poisson"),
        ("puce", "2. Observer ratio Déviance/ddl"),
        ("puce", "3. Basculer sur Binomial négatif => comparer AIC"),
        ("puce", "4. Calculer exp(beta_habitat) pour l'effet multiplicatif"),
        ("code", "exp(0.85) = 2.34  =>  abondance 2.34x plus élevée dans l'habitat de référence"),
        ("note", "Q2.1  Ratio = 3.4 pour Poisson. Quelle décision prendre ?\n"
                 "Q2.2  beta_haie = 0.62. Calculer et interpréter l'effet de la présence de haies."),

        ("h1", "Partie 3 - Modèles mixtes LMM (45 min)"),
        ("corps", "Contexte : longueur de tarse de poussins Mésange bleue dans 8 nichoirs (5 poussins/nichoir)."),
        ("puce", "1. Visualiser la variabilité inter-nichoirs vs intra-nichoir"),
        ("puce", "2. Comparer régression simple vs LMM (nichoir = effet aléatoire)"),
        ("puce", "3. Observer les erreurs standards des coefficients fixes"),
        ("note", "Q3.1  ICC = 0.55. Que signifie ce chiffre ?\n"
                 "Q3.2  L'erreur standard est plus grande en LMM. Pourquoi est-ce le comportement attendu ?"),

        ("h1", "Partie 4 - Domaines vitaux MCP et KDE (45 min)"),
        ("corps", "Contexte : 120 positions GPS d'un Busard des roseaux mâle en période de reproduction."),
        ("puce", "MCP : calculer MCP 100% puis 95%, comparer les surfaces"),
        ("puce", "KDE : observer isopleths 50% et 95%, faire varier la bande passante h"),
        ("note", "Q4.1  MCP 100% = 1850 ha, MCP 95% = 620 ha. Que révèle cet écart ?\n"
                 "Q4.2  Que se passe-t-il sur la forme du domaine vital quand h est très grand ?\n"
                 "Q4.3  MCP ou KDE pour identifier les zones de nidification ?"),
        ])


def pdf_tp3():
    _build_pdf("07_TP_Seance3.pdf",
        "TP Séance 3 - Dynamique des populations (fondamentaux)",
        "Durée 3h  -  Prérequis : Séances 1-2  -  Modules : Diversité, Croissance, Leslie, CMR",
        [
        ("h1", "Partie 1 - Richesse & Diversité (30 min)"),
        ("corps", "Comparer roselière (forte dominance), prairie humide (équilibrée), boisement alluvial."),
        ("puce", "1. Calculer H', Simpson (1-D), J pour chaque communauté"),
        ("puce", "2. Tracer les courbes d'accumulation - observer l'asymptote"),
        ("note", "Q1.1  Le boisement a S = 18 mais est-il le plus divers selon H' ?\n"
                 "Q1.2  Quelle communauté a J proche de 1 ? Que signifie cela biologiquement ?"),

        ("h1", "Partie 2 - Croissance exponentielle et logistique (45 min)"),
        ("corps", "Contexte : Milan royal réintroduit - N0 = 10 couples, r = 0.12, K = 80, horizon = 30 ans."),
        ("puce", "1. Observer la courbe exponentielle => valeur à 30 ans"),
        ("puce", "2. Activer la courbe logistique => stabilisation vers K"),
        ("puce", "3. Réduire K de 80 à 40 (perte d'habitat)"),
        ("puce", "4. Régler r = -0.08 (déclin) à partir de N0 = 60"),
        ("code", "T_1/2 = ln(2) / |r|    (temps de demi-vie pour r négatif)"),
        ("note", "Q2.1  À l'an 30, effectif exponentiel ? Est-ce réaliste ?\n"
                 "Q2.2  À quel moment la croissance nette est-elle maximale ?\n"
                 "Q2.3  r = -0.08 : temps de demi-vie = ?"),

        ("h1", "Partie 3 - Matrices de Leslie (60 min)"),
        ("corps", "Espèce : Cigogne blanche - 3 classes d'âge (juvénile, subadulte, adulte)."),
        ("tableau",
            [["Paramètre", "Valeur", "Signification"],
             ["F0", "0", "Juvéniles ne se reproduisent pas"],
             ["F1", "0.3", "Jeunes femelles par subadulte"],
             ["F2", "0.9", "Jeunes femelles par adulte"],
             ["S0", "0.45", "Survie juvénile"],
             ["S1", "0.72", "Survie subadulte"],
             ["S2", "0.85", "Survie adulte"],
            ],
            [3.0*cm, 2.5*cm, 10.0*cm]),
        ("puce", "1. Entrer les paramètres => calculer lambda"),
        ("puce", "2. Projeter sur 30 ans => observer la structure stable des âges"),
        ("puce", "3. Lire les élasticités => identifier le paramètre à prioriser"),
        ("puce", "4. Scénario A : S0 => 0.60 | Scénario B : S2 => 0.90 => comparer les lambda"),
        ("note", "Q3.2  lambda = ? => croissance, stabilité ou déclin ?\n"
                 "Q3.4  Élasticité S2 = 0.52 vs F2 = 0.18 => quelle intervention est prioritaire ?\n"
                 "Q3.5  S2 => 0.90 vs S0 => 0.60 : laquelle améliore le plus lambda ?"),

        ("h1", "Partie 4 - CMR (30 min)"),
        ("corps", "Contexte : Bruant des roseaux - Session 1 : M = 60 bagués. "
                  "Session 2 : C = 45 captures, R = 18 bagues recapturés."),
        ("code", "N_estimé = M x C / R = 60 x 45 / 18 = 150 individus"),
        ("puce", "1. Vérifier ce calcul dans ORNI-LAB => lire l'IC 95%"),
        ("puce", "2. Faire varier R de 5 à 40 => observer N_estimé et l'IC"),
        ("note", "Q4.1  L'IC est-il large ou étroit ? Que cela implique-t-il pour la gestion ?\n"
                 "Q4.2  Des bagues tombent entre les sessions => biais attendu sur N_estimé ?"),
        ])


def pdf_tp4():
    _build_pdf("08_TP_Seance4.pdf",
        "TP Séance 4 - Dynamique avancée et conservation",
        "Durée 3h  -  Prérequis : Séances 1-3  -  Modules : Occupation, Distance, Mann-Kendall, PVA, Scénarios",
        [
        ("h1", "Partie 1 - Modèles d'occupation (30 min)"),
        ("corps", "Contexte : Pic noir sur 50 placettes, 4 visites répétées, psi_vrai = 0.70, p = 0.35."),
        ("puce", "1. Observer occupation naïve = 0.25 (biais = 64% !)"),
        ("puce", "2. Ajuster le modèle MacKenzie => lire psi_estimé et p_estimé"),
        ("puce", "3. Réduire à 2 visites vs 6 visites => comparer la précision"),
        ("note", "Q1.1  Calculer le biais de l'occupation naïve.\n"
                 "Q1.4  Suivi 5 ans sans correction : occupation 62=>48%. Corrigé : psi stable mais p diminue. Interpréter."),

        ("h1", "Partie 2 - Distance sampling (30 min)"),
        ("corps", "Contexte : 180 détections de Tourterelle des bois sur 15 km de transects, distances 0-150 m."),
        ("puce", "1. Ajuster half-normal, puis hazard-rate => comparer AIC"),
        ("puce", "2. Lire D_estimé (individus/km²) et le CV"),
        ("note", "Q2.2  Si la zone fait 150 km² => abondance estimée ?\n"
                 "Q2.3  CV = 18% : précision acceptable pour conservation ?"),

        ("h1", "Partie 3 - Séries temporelles (30 min)"),
        ("corps", "Contexte : Gobemouche gris, 25 ans de données STOC, diminution apparente."),
        ("puce", "1. Appliquer Mann-Kendall => lire tau et p-value"),
        ("puce", "2. Lire la pente de Sen (individus perdus/an)"),
        ("puce", "3. Augmenter la variabilité interannuelle => effet sur la signification"),
        ("note", "Q3.1  tau = -0.38, p = 0.01 => conclusion ?\n"
                 "Q3.2  Pente = -3.2 ind/an, N0 = 85 => demi-vie = ?"),

        ("h1", "Partie 4 - PVA et Scénarios (60 min)"),
        ("corps", "Contexte : Vautour moine en réintroduction."),
        ("tableau",
            [["Paramètre", "Valeur"],
             ["N0", "35 individus reproducteurs"],
             ["r-moyen", "0.04 / an"],
             ["sigma", "0.18"],
             ["K", "200"],
             ["Seuil quasi-extinction", "15 individus"],
             ["Pertes annuelles", "3 (collisions lignes électriques)"],
             ["Horizon", "50 ans"],
            ],
            [5.0*cm, 10.5*cm]),
        ("puce", "1. Scénario de référence => risque et médiane finale"),
        ("puce", "2. Augmenter sigma à 0.30 => observer l'effet (r reste positif !)"),
        ("puce", "3. Éliminer les collisions (pertes => 0) => nouveau risque"),
        ("puce", "4. Renforcement : N0 => 70 => comparer avec élimination des collisions"),
        ("puce", "5. Module Scénarios : comparer A (réf), B (pertes=0), C (K=350), D (N0=70 + pertes=0)"),
        ("note", "Q4.2  r = 0.04 mais risque augmente avec sigma. Pourquoi r positif ne garantit-il pas la persistance ?\n"
                 "Q4.5  Classer les scénarios par efficacité et justifier votre recommandation.\n"
                 "Q4.6  Paramètres estimés sur 8 ans seulement. Comment adapter votre discours aux décideurs ?"),

        ("h1", "Récapitulatif - Les 18 modules en un coup d'oeil"),
        ("tableau",
            [["Module", "Question biologique", "Donnée minimale"],
             ["Analyse CSV", "Explorer mes données de terrain", "Fichier CSV"],
             ["Statistiques descriptives", "Résumer une variable", "Simulation"],
             ["Corrélation & régression", "Deux variables liées ?", "Simulation"],
             ["Tests statistiques", "Comparer 2-3 groupes", "Simulation"],
             ["GLM Comptage", "Modéliser des abondances", "Simulation"],
             ["Modèles mixtes", "Données groupées", "Simulation"],
             ["Domaine vital MCP", "Polygone englobant GPS", "Coordonnées GPS"],
             ["Domaine vital KDE", "Domaine vital probabiliste", "Coordonnées GPS"],
             ["Richesse & Diversité", "Comparer des communautés", "Abondances par espèce"],
             ["Croissance", "Projeter N dans le temps", "Simulation"],
             ["Matrices de Leslie", "Quelle classe d'âge protéger ?", "Fécondités + survies"],
             ["CMR", "Estimer une population cachée", "M, C, R"],
             ["Modèles d'occupation", "Corriger la détectabilité", "Visites répétées"],
             ["Distance sampling", "Densité sur transects", "Distances de détection"],
             ["Lotka-Volterra", "Dynamique proie-prédateur", "Simulation"],
             ["Séries temporelles", "Tendance sur plusieurs années", "CSV année + comptage"],
             ["PVA & Conservation", "Risque d'extinction à 50 ans", "N0, r, sigma, K, seuil"],
             ["Scénarios de gestion", "Quelle action prioriser ?", "Paramètres PVA"],
            ],
            [5.0*cm, 6.0*cm, 4.5*cm]),
        ])


# ══════════════════════════════════════════════════════════════════════════
#  POINT D'ENTRÉE
# ══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("\n=== Génération des supports de formation ORNI-LAB ===\n")
    print("PowerPoint :")
    generer_pptx()
    print("\nPDFs :")
    pdf_programme()
    pdf_guide_installation()
    pdf_fiches_biostat()
    pdf_fiches_dynamique()
    pdf_tp1()
    pdf_tp2()
    pdf_tp3()
    pdf_tp4()
    print(f"\nOK  Tous les fichiers sont dans :  {OUT_DIR}\n")
